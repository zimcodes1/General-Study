"""
File extraction utilities for different file types.
Handles PDF, DOCX, PPTX, Images (OCR), and plain text.
"""

import os
import logging
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

logger = logging.getLogger(__name__)


class TextExtractor(ABC):
    """Base class for text extraction."""

    @abstractmethod
    def extract(self, file_path: str) -> str:
        """Extract text from file. Returns text string."""
        pass


class PDFExtractor(TextExtractor):
    """Extract text from PDF files using PyMuPDF with pdfplumber fallback."""

    def extract(self, file_path: str) -> str:
        """Extract text from PDF using PyMuPDF, fallback to pdfplumber."""
        try:
            # Try PyMuPDF first (faster and more reliable)
            text = self._extract_with_pymupdf(file_path)
            if text.strip():
                logger.info(f"Successfully extracted {len(text)} chars from PDF using PyMuPDF")
                return text
            else:
                logger.warning("PyMuPDF returned empty text, trying pdfplumber")
        except Exception as e:
            logger.warning(f"PyMuPDF failed: {e}, trying pdfplumber")

        try:
            # Fallback to pdfplumber
            text = self._extract_with_pdfplumber(file_path)
            logger.info(f"Successfully extracted {len(text)} chars from PDF using pdfplumber")
            return text
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            raise

    @staticmethod
    def _extract_with_pymupdf(file_path: str) -> str:
        """Extract text using PyMuPDF (fitz)."""
        text = []
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc):
            page_text = page.get_text()
            text.append(page_text)
        doc.close()
        return "\n".join(text)

    @staticmethod
    def _extract_with_pdfplumber(file_path: str) -> str:
        """Extract text using pdfplumber."""
        text = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return "\n".join(text)


class DocxExtractor(TextExtractor):
    """Extract text from DOCX files."""

    def extract(self, file_path: str) -> str:
        """Extract text from DOCX."""
        try:
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    text.append(" | ".join(row_text))
            result = "\n".join(text)
            logger.info(f"Successfully extracted {len(result)} chars from DOCX")
            return result
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise


class PptxExtractor(TextExtractor):
    """Extract text from PPTX files."""

    def extract(self, file_path: str) -> str:
        """Extract text from PPTX."""
        try:
            presentation = Presentation(file_path)
            text = []
            for slide_num, slide in enumerate(presentation.slides):
                text.append(f"\n--- Slide {slide_num + 1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            result = "\n".join(text)
            logger.info(f"Successfully extracted {len(result)} chars from PPTX")
            return result
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise


class ImageExtractor(TextExtractor):
    """Extract text from images using OCR (Tesseract)."""

    def extract(self, file_path: str) -> str:
        """Extract text from image using OCR."""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            logger.info(f"Successfully extracted {len(text)} chars from image using OCR")
            return text
        except Exception as e:
            logger.error(f"Image OCR extraction failed: {e}")
            raise


class TextFileExtractor(TextExtractor):
    """Extract (read) plain text files."""

    def extract(self, file_path: str) -> str:
        """Read plain text file."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            logger.info(f"Successfully read {len(text)} chars from text file")
            return text
        except Exception as e:
            logger.error(f"Text file reading failed: {e}")
            raise


class ExtractorFactory:
    """Factory for selecting appropriate text extractor based on file type."""

    EXTRACTORS = {
        "pdf": PDFExtractor(),
        "doc": DocxExtractor(),
        "docx": DocxExtractor(),
        "ppt": PptxExtractor(),
        "pptx": PptxExtractor(),
        "image": ImageExtractor(),
        "jpg": ImageExtractor(),
        "jpeg": ImageExtractor(),
        "png": ImageExtractor(),
        "gif": ImageExtractor(),
        "webp": ImageExtractor(),
        "txt": TextFileExtractor(),
    }

    @classmethod
    def get_extractor(cls, file_type: str) -> TextExtractor:
        """Get appropriate extractor for file type."""
        extractor = cls.EXTRACTORS.get(file_type.lower())
        if not extractor:
            raise ValueError(f"Unsupported file type: {file_type}")
        return extractor


def _infer_type_from_path(file_path: str) -> str:
    ext = Path(file_path).suffix.lower().lstrip(".")
    if not ext:
        return "other"
    if ext in {"jpg", "jpeg", "png", "gif", "webp"}:
        return "image"
    return ext


def _normalize_file_type(file_path: str, file_type: str) -> str:
    if not file_type:
        return _infer_type_from_path(file_path)
    normalized = file_type.lower()
    if normalized in {"document", "other"}:
        return _infer_type_from_path(file_path)
    return normalized


def extract_text_from_file(file_path: str, file_type: str) -> str:
    """
    Extract text from a file based on its type.

    Args:
        file_path: Full path to the file
        file_type: File type/extension (pdf, docx, pptx, image, etc.)

    Returns:
        Extracted text as string

    Raises:
        ValueError: If file type is unsupported
        Exception: If extraction fails
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    normalized_type = _normalize_file_type(file_path, file_type)
    extractor = ExtractorFactory.get_extractor(normalized_type)
    return extractor.extract(file_path)
